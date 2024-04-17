import os
import shutil
import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.llms import OpenAI
import pandas as pd
from pptx import Presentation
from docx import Document
from dotenv import load_dotenv
from time import sleep
from requests.exceptions import HTTPError
 
load_dotenv()
 
class SessionState:
    def __init__(self, **kwargs):
        self.__dict__.update(kwargs)
 
def process_documents(documents_folder):
    documents = []
    for file_name in os.listdir(documents_folder):
        file_path = os.path.join(documents_folder, file_name)
        if file_path.endswith(".pdf"):
            pdf_reader = PdfReader(file_path)
            for page in pdf_reader.pages:
                documents.append(page.extract_text())
        elif file_path.endswith(".txt"):
            with open(file_path, "r") as file:
                documents.append(file.read())
        elif file_path.endswith(".docx"):
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                documents.append(paragraph.text)
        elif file_path.endswith(".xlsx"):
            xls = pd.ExcelFile(file_path)
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                for col in df.columns:
                    documents.extend(df[col].dropna().astype(str).tolist())
        elif file_path.endswith(".pptx"):
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        documents.append(shape.text)
 
    text = "\n".join(documents)
    if not text:
        return None
 
    char_text_splitter = CharacterTextSplitter(separator="\n", chunk_size=1000, chunk_overlap=200, length_function=len)
    text_chunks = char_text_splitter.split_text(text)
    if not text_chunks:
        return None
 
    embeddings = OpenAIEmbeddings()
    docsearch = FAISS.from_texts(text_chunks, embeddings)
    return docsearch
 
def answer_query(query, docsearch, chain):
    docs = docsearch.similarity_search(query)
    response = chain.run(input_documents=docs, question=query)
    return response
 
def delete_files(selected_documents, documents_folder):
    for file_name in selected_documents:
        file_path = os.path.join(documents_folder, file_name)
        try:
            os.remove(file_path)
            st.write(f"{file_name} deleted successfully.")
        except Exception as e:
            st.error(f"Error deleting {file_name}: {e}")
 
def exponential_backoff_retry(func, max_retries=5, initial_wait_time=1):
    retries = 0
    wait_time = initial_wait_time
    while retries < max_retries:
        try:
            return func()
        except HTTPError as e:
            if e.response.status_code == 429:  # Rate limit exceeded
                st.warning(f"Rate limit exceeded. Retrying in {wait_time} seconds...")
                sleep(wait_time)
                wait_time *= 2  # Exponential backoff
                retries += 1
            else:
                raise  # Reraise the exception if it's not a rate limit error
    st.error("Max retry attempts reached. Please try again later.")
 
def display_assistant(session_state, chain):
    # AI assistant chat interface
    st.sidebar.header("Assistant Bot🤖")
    chat_history = st.sidebar.empty()  # Placeholder for chat history
    user_input = st.sidebar.text_input("Query:", key="user_input")
    if st.sidebar.button("Submit"):
        response = answer_query(user_input, session_state.docsearch, chain)
        st.sidebar.text("Bot Response:")
        #st.sidebar.text(response)
        # Concatenate the lines of the response into a single string
        horizontal_response = " ".join(response.split('\n'))
        
        # Display the horizontally formatted response
        st.sidebar.write(horizontal_response)
 
def main():
    st.header("Data Explorer Bot🔍")
    session_state = SessionState(documents_folder=None, docsearch=None, response="", is_admin=False)
 
    user_type = st.radio("Choose User Role:", ["Admin", "Normal User"])
 
    if user_type == "Admin":
        admin_username = st.text_input("Enter admin username:")
        admin_password = st.text_input("Enter admin password:", type="password")
        login_button = st.button("Login")
 
        if admin_username and admin_password:
            if admin_username == "admin" and admin_password == "admin@123":
                session_state.is_admin = True
                st.success("Admin access granted!")
 
                documents_folder_admin = "uploaded_documents_admin"
                if not os.path.exists(documents_folder_admin):
                    os.makedirs(documents_folder_admin)
 
                uploaded_files = st.file_uploader("Upload your documents", accept_multiple_files=True,
                                                  type=["pdf", "txt", "docx", "xlsx", "pptx"])
                for file in uploaded_files:
                    with open(os.path.join(documents_folder_admin, file.name), "wb") as f:
                        f.write(file.getvalue())
 
                session_state.documents_folder = documents_folder_admin
                session_state.docsearch = process_documents(documents_folder_admin)
                openai_api_key = os.getenv("OPENAI_API_KEY")
 
                def load_chain():
                    return load_qa_chain(OpenAI(api_key=openai_api_key), chain_type="stuff")
 
                chain = exponential_backoff_retry(load_chain)
 
                st.subheader("Uploaded Documents📑")
                selected_documents = []
                for file_name in os.listdir(documents_folder_admin):
                    selected = st.checkbox(file_name)
                    if selected:
                        selected_documents.append(file_name)
 
                if st.button("Delete Documents"):
                    delete_files(selected_documents, documents_folder_admin)
 
                display_assistant(session_state, chain)
 
    else:  # Normal user functionality
        documents_folder_admin = "uploaded_documents_admin"
        documents_folder_normal_user = "uploaded_Documents"
 
        if not os.path.exists(documents_folder_admin):
            st.write("No documents uploaded yet by the admin. Please wait for the admin to upload documents.")
            return
 
        if not os.path.exists(documents_folder_normal_user):
            os.makedirs(documents_folder_normal_user)
 
        # Copy documents from admin folder to normal user folder
        for file_name in os.listdir(documents_folder_admin):
            src_path = os.path.join(documents_folder_admin, file_name)
            dest_path = os.path.join(documents_folder_normal_user, file_name)
            shutil.copy(src_path, dest_path)
 
        session_state.documents_folder = documents_folder_normal_user
        session_state.docsearch = process_documents(documents_folder_normal_user)
        openai_api_key = os.getenv("OPENAI_API_KEY")
 
        def load_chain():
            return load_qa_chain(OpenAI(api_key=openai_api_key), chain_type="stuff")
 
        chain = exponential_backoff_retry(load_chain)
 
        st.subheader("Uploaded Documents📑")
        for file_name in os.listdir(documents_folder_normal_user):
            st.write(file_name)
           
        display_assistant(session_state, chain)
 
if __name__ == "__main__":
    main()